#!/usr/bin/env python
# Builds the Track & Trace showcase deck from the Blue Ridge kickoff template,
# inheriting its branded master (border frame, corner logo, palette). Keeps the
# Double Logo title/close, strips the kickoff slides, adds feature slides.
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

TEMPLATE = r"C:\Users\JACQUE~1\AppData\Local\Temp\claude\C--Users-JacquesPotgieter-Documents-Dev-MPP\d661cf5e-71a5-47c1-b545-68dae760b4a0\scratchpad\template.pptx"
OUT = r"C:\Users\JacquesPotgieter\Documents\Dev\MPP\docs\showcase\Track-and-Trace-Showcase.pptx"
IMG = r"C:\Users\JacquesPotgieter\Documents\Dev\MPP\docs\showcase\screenshots"

NAVY = RGBColor(0x1E, 0x27, 0x61)
CYAN = RGBColor(0x29, 0xAB, 0xE2)
GOLD = RGBColor(0xF7, 0xA8, 0x23)
GRAY = RGBColor(0x5A, 0x64, 0x72)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TITLE_FONT = "Cambria"
BODY_FONT = "Calibri"

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation(TEMPLATE)
layouts = {l.name: l for l in prs.slide_masters[0].slide_layouts}

def layout(name):
    return layouts[name]

# ---- strip kickoff slides, keep slide 1 (title). Drop the relationship too so
# the orphaned slide part is not written (else duplicate partnames corrupt the file). ----
sldIdLst = prs.slides._sldIdLst
ids = list(sldIdLst)
for sid in ids[1:]:
    rId = sid.get(qn('r:id'))
    try:
        prs.part.drop_rel(rId)
    except Exception:
        pass
    sldIdLst.remove(sid)

def add_blank():
    return prs.slides.add_slide(layout("Blank"))

def textbox(slide, l, t, w, h, text, size, color, font=BODY_FONT, bold=False,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False, spacing=None):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing is not None: p.line_spacing = spacing
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.name = font; r.font.color.rgb = color
    return tb

def title_caption(slide, title, caption):
    textbox(slide, Inches(0.62), Inches(0.4), Inches(12.1), Inches(0.75),
            title, 30, NAVY, font=TITLE_FONT, bold=True)
    if caption:
        textbox(slide, Inches(0.64), Inches(1.16), Inches(12.06), Inches(0.5),
                caption, 15, GRAY, font=BODY_FONT)

def img_dims(path):
    from PIL import Image
    with Image.open(path) as im:
        return im.size

def place_image_fit(slide, path, box_l, box_t, box_w, box_h):
    iw, ih = img_dims(path)
    ar = iw / ih
    bar = box_w / box_h
    if ar > bar:
        w = box_w; h = int(box_w / ar)
    else:
        h = box_h; w = int(box_h * ar)
    l = int(box_l + (box_w - w) / 2)
    t = int(box_t + (box_h - h) / 2)
    pic = slide.shapes.add_picture(path, l, t, width=w, height=h)
    # subtle line border
    pic.line.color.rgb = RGBColor(0xD9, 0xDD, 0xE3); pic.line.width = Pt(0.75)
    return pic

def single(title, caption, img):
    s = add_blank()
    title_caption(s, title, caption)
    place_image_fit(s, f"{IMG}\\{img}", Inches(0.7), Inches(1.75), Inches(11.93), Inches(5.4))
    return s

def multi(title, caption, imgs, labels=None):
    s = add_blank()
    title_caption(s, title, caption)
    n = len(imgs); gap = Inches(0.3)
    total_w = Inches(12.33)
    iw = int((total_w - gap * (n - 1)) / n)
    box_h = Inches(4.6) if n == 2 else Inches(3.2)
    top = Inches(2.05) if n == 2 else Inches(2.5)
    x = Inches(0.5)
    for i, im in enumerate(imgs):
        place_image_fit(s, f"{IMG}\\{im}", x, top, iw, box_h)
        if labels:
            textbox(s, x, top + box_h + Inches(0.08), iw, Inches(0.35),
                    labels[i], 12, GRAY, align=PP_ALIGN.CENTER)
        x += iw + gap
    return s

def section(title, sub):
    s = prs.slides.add_slide(layout("Section Header"))
    # fill title/body placeholders if present, else add textboxes
    filled_title = False
    for ph in s.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = title
            for p in ph.text_frame.paragraphs:
                for r in p.runs:
                    r.font.color.rgb = NAVY; r.font.name = TITLE_FONT
            filled_title = True
        elif ph.placeholder_format.idx == 1:
            ph.text = sub
            for p in ph.text_frame.paragraphs:
                for r in p.runs:
                    r.font.color.rgb = GRAY; r.font.name = BODY_FONT
    if not filled_title:
        textbox(s, Inches(0.9), Inches(2.9), Inches(11.5), Inches(1.2), title, 40, NAVY,
                font=TITLE_FONT, bold=True)
        textbox(s, Inches(0.92), Inches(4.15), Inches(11.5), Inches(0.8), sub, 17, GRAY)
    return s

# ---- Blank the "Kickoff Meeting" prompt text on the Double Logo layout (used by
# both the title and closing slides) so it doesn't show through. ----
for shp in layout("Double Logo").shapes:
    if shp.has_text_frame and "kickoff" in shp.text_frame.text.lower():
        for p in shp.text_frame.paragraphs:
            for r in list(p.runs):
                r.text = ""

# ---- Slide 1: title text, centered under the logo ----
s1 = prs.slides[0]
textbox(s1, Inches(0.9), Inches(5.45), Inches(11.53), Inches(0.8),
        "Track & Trace + Genealogy", 34, NAVY, font=TITLE_FONT, bold=True, align=PP_ALIGN.CENTER)
textbox(s1, Inches(0.9), Inches(6.35), Inches(11.53), Inches(0.5),
        "Manufacturing Execution System  \u2014  Capability Showcase", 16, GRAY,
        font=BODY_FONT, align=PP_ALIGN.CENTER)

# ---- Track & Trace ----
section("Track & Trace + Genealogy", "Complete part lineage \u2014 every LOT, from raw casting to shipment.")
single("Full LOT traceability",
       "Every LOT carries its complete production history, genealogy, as-built BOM, and quality record \u2014 one screen.",
       "C01_lot_detail.png")
single("Genealogy \u2014 forward & backward lineage",
       "Trace any finished good to the exact castings and components it consumed, and any lot forward to what it became.",
       "C02_genealogy.png")
single("Find any part \u2014 live location & status",
       "Search any LOT, vendor lot, or part number and see where every unit sits on the floor and its current state.",
       "C03_lot_search.png")

# ---- Reporting ----
section("Reporting", "On-demand operational and traceability reporting, straight from the floor.")
single("Plant-wide inventory, on demand",
       "A live WIP snapshot by item and location \u2014 total pieces, open LOTs, and where everything is, as of this minute.",
       "C07_report_inventory.png")
multi("Operational reporting suite",
      "Tooling life and line performance \u2014 shots per die vs. limit, and weekly output, scrap, and downtime by line.",
      ["C08_report_shotcount.png", "C09_report_lineperf.png"],
      labels=["Die Cast Shot Count \u2014 die life vs. limit", "Production Line Performance \u2014 weekly by line"])

# ---- Configuration ----
section("Configuration", "The entire system is operator-configurable \u2014 no code.")
multi("Configure the plant, parts, and processes",
      "Model the plant hierarchy, define parts with routes / BOMs / container rules, and version every operation template.",
      ["C10_config_plant.png", "C11_config_items.png", "C12_config_optemplates.png"],
      labels=["Plant hierarchy (ISA-95)", "Item master + routes / BOMs", "Versioned operation templates"])

# ---- Compliance capstone ----
sc = section("Built on a Data-Integrity Foundation",
             "Architected to support the record-keeping and traceability requirements of regulated manufacturing.")

# Compliance content slide: table (left) + audit image (right)
sc2 = add_blank()
title_caption(sc2, "Every change, attributable and traceable",
              "A complete, append-only audit trail \u2014 who changed what, when, with before / after values.")
rows = [
    ("Append-only event log + full audit trail (who / what / when, before\u2192after)",
     "21 CFR Part 11 audit trail \u00b7 ALCOA+"),
    ("AD identity + audited per-action elevation (grant and deny recorded)",
     "Part 11 electronic-signature-style authorization"),
    ("No hard deletes \u2014 soft-delete + UTC-stamped, immutable records",
     "Record protection & retention"),
    ("Full lot genealogy + production history + as-built BOM",
     "Device History Record traceability (ISO 13485 / 820)"),
    ("Spec-driven quality + hold / nonconformance control",
     "Nonconformance handling"),
    ("Draft \u2192 Published \u2192 Deprecated version control on BOMs, routes, specs",
     "Controlled specification revision control"),
]
# left column: mapping list
ly = Inches(1.9)
for feat, std in rows:
    textbox(sc2, Inches(0.7), ly, Inches(6.2), Inches(0.55), feat, 12.5, NAVY, bold=True)
    textbox(sc2, Inches(0.7), ly + Inches(0.42), Inches(6.2), Inches(0.4), std, 11.5, CYAN, italic=True)
    ly += Inches(0.86)
# right: audit screenshot
place_image_fit(sc2, f"{IMG}\\C13_config_audit.png", Inches(7.15), Inches(1.9), Inches(5.55), Inches(4.9))
textbox(sc2, Inches(7.15), Inches(6.85), Inches(5.55), Inches(0.4),
        "Live configuration audit log \u2014 attribution, change diffs, severity.", 11, GRAY, align=PP_ALIGN.CENTER)

# ---- Close: Double Logo ----
close = prs.slides.add_slide(layout("Double Logo"))
textbox(close, Inches(0.9), Inches(5.7), Inches(11.53), Inches(0.7),
        "Track & Trace + Genealogy", 26, NAVY, font=TITLE_FONT, bold=True, align=PP_ALIGN.CENTER)
textbox(close, Inches(0.9), Inches(6.45), Inches(11.53), Inches(0.5),
        "Transforming Manufacturing Beyond Limits", 15, GRAY, font=BODY_FONT, italic=True,
        align=PP_ALIGN.CENTER)

# ---- Speaker notes (presenter talk track) ----
NOTES = [
    "Opening. This is a working Manufacturing Execution System for an aluminum die-casting operation supporting automotive. Everything shown is live from the running system. Three themes: complete track & trace, on-demand reporting, and a data-integrity foundation for regulated manufacturing.",
    "Track & Trace section. The core promise: every part carries its full lineage from raw casting through machining, assembly, and shipment.",
    "This is a single finished-good LOT. One screen shows its whole life: every movement, status change, the components it consumed, its as-built BOM version, and its quality record. No paper, no spreadsheets.",
    "Genealogy works both directions: from a finished good down to the exact castings and components consumed, and from any raw lot forward to everything it became. This is what a full recall or containment investigation needs.",
    "Operators or quality find any lot instantly by lot number, vendor lot, or part number, and immediately see where every unit is on the floor and its current state.",
    "Reporting section. Every report is on-demand from the floor and exports to PDF.",
    "Current inventory: a live plant-wide WIP snapshot by item and location. Total pieces, open lots, and exactly where everything sits, as of this minute.",
    "Operational reports. Left: tooling life, shots per die against its limit with near/over-limit flags. Right: weekly production, scrap, and downtime by line for trend visibility.",
    "Configuration section. The whole system is configured by engineers through the UI. No code, no vendor change orders.",
    "Model the plant hierarchy in ISA-95 tiers, define every part with its routes, BOMs, and container rules, and version every operation template. This is why the system adapts to a new part or line without a project.",
    "Compliance section. This matters most for regulated manufacturing. The point is not a certificate; it is that the data model is architected to support these requirements.",
    "Every capability maps to a real requirement. Append-only audit trail with before/after values and attribution supports 21 CFR Part 11 and ALCOA+. AD-backed per-action elevation, recorded on grant and deny, is the electronic-signature-style control. No hard deletes, full genealogy, spec-driven quality, and Draft-Published-Deprecated version control round it out. On the right is the live audit log itself.",
    "Close. To recap: complete genealogy, live floor visibility, instant reporting, full configurability, and a compliance-ready record foundation. Happy to go deeper on any screen.",
]
for slide, note in zip(prs.slides, NOTES):
    if note:
        slide.notes_slide.notes_text_frame.text = note

prs.save(OUT)
print("saved", OUT, "slides:", len(prs.slides._sldIdLst))
